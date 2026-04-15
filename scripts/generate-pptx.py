#!/usr/bin/env python3
"""
Generate PowerPoint from JSON content using Nationwide template.

Usage:
    python scripts/generate-pptx.py content.json output.pptx
    echo '{"title": "...", "slides": [...]}' | python scripts/generate-pptx.py - output.pptx
"""

import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("Error: python-pptx not installed. Run: pip install python-pptx")
    sys.exit(1)


def load_template():
    """Load the Nationwide template."""
    template_path = Path(__file__).parent.parent / "templates/pct/nationwide_default.pptx"
    if not template_path.exists():
        print(f"Error: Template not found at {template_path}")
        sys.exit(1)
    return Presentation(template_path)


def add_title_slide(prs, slide_data):
    """Add a title slide."""
    layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(layout)

    title = slide.shapes.title
    subtitle = slide.placeholders[1] if len(slide.placeholders) > 1 else None

    title.text = slide_data.get("title", "")
    if subtitle and slide_data.get("subtitle"):
        subtitle.text = slide_data["subtitle"]

    return slide


def add_content_slide(prs, slide_data):
    """Add a content slide with bullets."""
    layout = prs.slide_layouts[1]  # Content layout
    slide = prs.slides.add_slide(layout)

    title = slide.shapes.title
    body = slide.placeholders[1] if len(slide.placeholders) > 1 else None

    title.text = slide_data.get("title", "")

    if body and slide_data.get("bullets"):
        tf = body.text_frame
        tf.clear()
        for i, bullet in enumerate(slide_data["bullets"]):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = bullet
            p.level = 0

    # Add speaker notes if provided
    if slide_data.get("notes"):
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = slide_data["notes"]

    return slide


def add_two_column_slide(prs, slide_data):
    """Add a two-column comparison slide."""
    layout = prs.slide_layouts[3]  # Two content layout (may vary by template)
    slide = prs.slides.add_slide(layout)

    title = slide.shapes.title
    title.text = slide_data.get("title", "")

    # Find the content placeholders
    placeholders = [p for p in slide.placeholders if p.placeholder_format.idx > 0]

    if len(placeholders) >= 2:
        left = placeholders[0]
        right = placeholders[1]

        # Left column
        if slide_data.get("left_bullets"):
            tf = left.text_frame
            tf.clear()
            if slide_data.get("left_title"):
                p = tf.paragraphs[0]
                p.text = slide_data["left_title"]
                p.font.bold = True
            for bullet in slide_data["left_bullets"]:
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 0

        # Right column
        if slide_data.get("right_bullets"):
            tf = right.text_frame
            tf.clear()
            if slide_data.get("right_title"):
                p = tf.paragraphs[0]
                p.text = slide_data["right_title"]
                p.font.bold = True
            for bullet in slide_data["right_bullets"]:
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 0

    return slide


def add_closing_slide(prs, slide_data):
    """Add a closing slide."""
    # Use content layout for closing with next steps
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)

    title = slide.shapes.title
    body = slide.placeholders[1] if len(slide.placeholders) > 1 else None

    title.text = slide_data.get("title", "Next Steps")

    if body and slide_data.get("bullets"):
        tf = body.text_frame
        tf.clear()
        for i, bullet in enumerate(slide_data["bullets"]):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = bullet
            p.level = 0

    return slide


def generate_presentation(content, output_path):
    """Generate a PowerPoint from content JSON."""
    prs = load_template()

    # Remove any existing slides from template
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]

    # Add slides based on content
    for slide_data in content.get("slides", []):
        slide_type = slide_data.get("type", "content")

        if slide_type == "title":
            add_title_slide(prs, slide_data)
        elif slide_type == "content":
            add_content_slide(prs, slide_data)
        elif slide_type == "two_column":
            add_two_column_slide(prs, slide_data)
        elif slide_type == "section":
            # Section slides are like title slides
            add_title_slide(prs, slide_data)
        elif slide_type == "closing":
            add_closing_slide(prs, slide_data)
        else:
            print(f"Warning: Unknown slide type '{slide_type}', using content layout")
            add_content_slide(prs, slide_data)

    # Save
    prs.save(output_path)
    print(f"Generated: {output_path}")
    print(f"Slides: {len(prs.slides)}")


def main():
    if len(sys.argv) < 3:
        print("Usage: python generate-pptx.py <content.json|-> <output.pptx>")
        print("  Use '-' to read JSON from stdin")
        sys.exit(1)

    input_arg = sys.argv[1]
    output_path = sys.argv[2]

    # Read content
    if input_arg == "-":
        content = json.load(sys.stdin)
    else:
        with open(input_arg) as f:
            content = json.load(f)

    generate_presentation(content, output_path)


if __name__ == "__main__":
    main()
